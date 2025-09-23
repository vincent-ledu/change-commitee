from __future__ import annotations
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN


def _parse_hex_color(value: str) -> RGBColor | None:
    if not isinstance(value, str):
        return None
    text = value.strip().lstrip('#')
    if len(text) != 6:
        return None
    try:
        r = int(text[0:2], 16)
        g = int(text[2:4], 16)
        b = int(text[4:6], 16)
    except ValueError:
        return None
    return RGBColor(r, g, b)


RFC_URL_TEMPLATE = "https://outils.change.fr/change={rfc}"


def set_rfc_url_template(template: str | None) -> None:
    """Override the base URL template for RFC hyperlinks. Must contain '{rfc}'."""
    global RFC_URL_TEMPLATE
    if template and isinstance(template, str) and "{rfc}" in template:
        RFC_URL_TEMPLATE = template


def hyperlink_for_rfc(rfc: str) -> str:
    try:
        return RFC_URL_TEMPLATE.format(rfc=str(rfc).lower())
    except Exception:
        return f"{RFC_URL_TEMPLATE}{str(rfc).lower()}"


_DEFAULT_COLOR_MAP = {
    "urgent": RGBColor(255, 140, 0),   # orange
    "normal": RGBColor(0, 102, 204),   # blue
    "agile": RGBColor(0, 153, 0),      # green
}


COLOR_MAP = _DEFAULT_COLOR_MAP.copy()
DEFAULT_COLOR = RGBColor(100, 100, 100)


def set_timeline_color_map(overrides: dict | None) -> None:
    """Override the timeline color map from config values like {'urgent': '#FF8800'}"""
    COLOR_MAP.clear()
    COLOR_MAP.update(_DEFAULT_COLOR_MAP)
    if not overrides or not isinstance(overrides, dict):
        return
    for key, value in overrides.items():
        norm_key = str(key or '').strip().lower()
        if not norm_key:
            continue
        parsed = _parse_hex_color(value)
        if parsed is None:
            print(f"[WARN] Ignoring invalid color for '{norm_key}': {value!r}")
            continue
        COLOR_MAP[norm_key] = parsed


def set_title(prs, slide, text: str, font_size_pt: int = 24) -> None:
    """Set a slide title into the layout's title placeholder if present.
    If no title placeholder, add a textbox at the top with a title style.
    """
    try:
        title_shape = slide.shapes.title
    except Exception:
        title_shape = None

    if title_shape is not None:
        try:
            tf = title_shape.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            run.font.size = Pt(font_size_pt)
            run.font.bold = True
            return
        except Exception:
            pass

    # No title placeholder; add a textbox at the top
    left = Cm(1.0)
    top = Cm(1.0)
    width = prs.slide_width - Cm(2.0)
    height = Cm(1.5)
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size_pt)
    run.font.bold = True

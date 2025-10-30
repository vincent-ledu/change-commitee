from __future__ import annotations

import pandas as pd
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

from layouts import choose_detail_layout
from .utils import set_title


def _extract_assignee_series(df: pd.DataFrame) -> pd.Series:
    """Collect a normalized assignee series using the best available column."""
    if df.empty:
        return pd.Series(dtype=object)

    candidates_primary = [
        "Affecté",
        "Affecte",
        "Assigné à",
        "Assignation",
        "Responsable",
        "Owner",
    ]
    candidates_fallback = [
        "Affecté à",
        "Groupe d’affectation",
        "Groupe gestionnaire",
    ]

    series = pd.Series([""] * len(df), index=df.index, dtype=object)

    def _merge(col: str) -> None:
        nonlocal series
        if col not in df.columns:
            return
        values = df[col].fillna("").astype(str).str.strip()
        mask = series.eq("") & values.ne("")
        if mask.any():
            series.loc[mask] = values.loc[mask]

    for col in candidates_primary:
        _merge(col)
    for col in candidates_fallback:
        _merge(col)

    return series


def add_assignee_bar_chart_slide(prs: Presentation,
                                 df: pd.DataFrame,
                                 title: str,
                                 layout_index: int | None = None) -> None:
    """Add a bar chart slide showing number of changes per assignee."""
    layout = choose_detail_layout(prs, layout_index)
    slide = prs.slides.add_slide(layout)
    set_title(prs, slide, title)

    if df.empty:
        top = Cm(4.0)
        left = Cm(2.0)
        width = prs.slide_width - Cm(4.0)
        height = Cm(2.0)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = "Aucun changement pour cette période."
        run.font.size = Pt(18)
        return

    series = _extract_assignee_series(df)
    if series.empty:
        counts = pd.Series([len(df)], index=["non affecté"])
    else:
        normalized = series.replace("", pd.NA).fillna("non affecté")
        normalized = normalized.astype(str).str.strip()
        normalized = normalized.mask(normalized.eq(""), "non affecté")
        counts = normalized.value_counts().sort_values(ascending=False)
        if counts.empty:
            counts = pd.Series([len(df)], index=["non affecté"])

    chart_data = ChartData()
    chart_data.categories = list(counts.index)
    chart_data.add_series("Changements", counts.tolist())

    # Chart placement: leave room for title above and optional notes below.
    left = Cm(1.5)
    top = Cm(3.5)
    width = prs.slide_width - Cm(3.0)
    height = prs.slide_height - Cm(6.0)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        left,
        top,
        width,
        height,
        chart_data,
    ).chart

    chart.has_title = False
    chart.has_legend = False
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.minimum_scale = 0
    chart.value_axis.tick_labels.font.size = Pt(12)
    chart.value_axis.format.line.color.rgb = RGBColor(191, 191, 191)
    chart.category_axis.tick_labels.font.size = Pt(12)
    chart.category_axis.tick_labels.orientation = 45
    chart.category_axis.format.line.color.rgb = RGBColor(191, 191, 191)

    # Apply a consistent fill color for the bars.
    series_obj = chart.series[0]
    fill_color = RGBColor(91, 155, 213)  # match default normal type color
    for point in series_obj.points:
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = fill_color

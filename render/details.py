from __future__ import annotations
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from layouts import choose_detail_layout
from .utils import hyperlink_for_rfc


DETAIL_FIELDS = [
    ("Description", "Description"),
    ("Justification", "Justification"),
    ("Plan d’implémentation", "Plan d’implémentation"),
    ("Analyse risques & impacts", "Analyse de risques et de l’impact"),
    ("Plan de retour arrière", "Plan de retour en arrière"),
    ("Plan de tests", "Plan de tests"),
    ("Informations complémentaires", "Informations complémentaires"),
]


def add_detail_slide(prs: Presentation, row: pd.Series, layout_index: int | None = None) -> None:
    layout_to_use = choose_detail_layout(prs, layout_index)

    slide = prs.slides.add_slide(layout_to_use)

    # Use title placeholder if present; otherwise fall back to a textbox
    title_shape = getattr(slide.shapes, "title", None)
    if title_shape is None:
        left = Cm(1.0); top = Cm(1.0); width = prs.slide_width - Cm(2.0); height = Cm(1.2)
        title_shape = slide.shapes.add_textbox(left, top, width, height)
    title_tf = title_shape.text_frame
    title_tf.clear()
    p = title_tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    rfc = str(row["Numéro"]).strip()
    resume = str(row.get("Résumé", "")).strip()

    run1 = p.add_run()
    run1.text = rfc
    run1.hyperlink.address = hyperlink_for_rfc(rfc)
    run1.font.size = Pt(28)
    run1.font.bold = True

    run2 = p.add_run()
    run2.text = f" — {resume}"
    run2.font.size = Pt(24)

    # Info badges (Type, État, Demandeur, Affecté, Affecté à, Début/Fin planifiées)
    # Rendered as small rounded rectangles (two per row: label, value) at the very top‑right.
    # We do NOT push the details table down; layout should leave room above the details table.
    badges_top = Cm(0.5)  # top margin for badges area (moved up by 0.5 cm)
    # Area reserved on the right for badges (fixed compact width)
    # Increased widths by 0.5 cm each as requested
    label_w = Cm(1.8 + 0.5)
    value_w = Cm(2.2 + 0.5)
    pair_w = label_w + value_w
    gap_h = Cm(0.06)
    x_label = prs.slide_width - Cm(1.0) - pair_w
    x_value = x_label + label_w

    # Prepare values (with robust date formatting)
    def _fmt_dt(col_dt: str, col_text: str) -> str:
        try:
            return row[col_dt].strftime("%d/%m/%Y %H:%M")
        except Exception:
            return str(row.get(col_text, "")).strip()

    badges = [
        ("Type", str(row.get("Type", "")).strip()),
        ("État", str(row.get("Etat", "")).strip()),
        ("Demandeur", str(row.get("Demandeur", "")).strip()),
        ("Affecté", str(row.get("Affecté", "")).strip()),
        ("Affecté à", str(row.get("Affecté à", "")).strip()),
        ("Début planifié", _fmt_dt("start_dt", "Date de début planifiée")),
        ("Fin planifiée", _fmt_dt("end_dt", "Date de fin planifiée")),
    ]

    # Determine badge colors based on change type (top-right cartouche)
    type_label = str(row.get("Type", "")).strip().lower()
    badge_palettes = {
        "normal": (
            RGBColor(91, 155, 213),  # blue label
            RGBColor(42, 96, 153),   # blue value
        ),
        "agile": (
            RGBColor(142, 197, 71),  # green label
            RGBColor(94, 155, 43),   # green value
        ),
        "urgent": (
            RGBColor(242, 171, 78),  # orange label
            RGBColor(209, 118, 6),   # orange value
        ),
    }
    color_label, color_value = badge_palettes.get(
        type_label,
        badge_palettes["normal"],  # default to blue when type unrecognized
    )

    # Draw rounded-rectangle badges if values present
    items = [(lab, val) for lab, val in badges if val]
    planned_tbl_top = Cm(4.6)
    if items:
        # Compute row height to fit before the details table top
        max_h = planned_tbl_top - badges_top - Cm(0.1)
        rows = len(items)
        row_h = max(Cm(0.38), min(Cm(0.6), (max_h - gap_h * (rows - 1)) / max(1, rows)))

        def _add_round_box(x, y, w, h, text, fill_rgb: RGBColor, bold=False):
            shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
            shp.fill.solid()
            shp.fill.fore_color.rgb = fill_rgb
            shp.line.fill.background()
            tf = shp.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            r = p.add_run()
            r.text = text
            r.font.size = Pt(6)
            r.font.bold = bool(bold)
            r.font.color.rgb = RGBColor(255, 255, 255)

        def _needs_state_alert(badge_label: str, badge_value: str) -> bool:
            if type_label != "normal" or badge_label.lower() != "état":
                return False
            import unicodedata
            norm = ''.join(
                c for c in unicodedata.normalize('NFD', str(badge_value or "").strip().lower())
                if unicodedata.category(c) != 'Mn'
            )
            return norm not in {"planifie", "autoriser"}

        for i, (lab, val) in enumerate(items):
            y = badges_top + i * (row_h + gap_h)
            label_color = color_label
            value_color = color_value
            if _needs_state_alert(lab, val):
                label_color = RGBColor(242, 171, 78)
                value_color = RGBColor(209, 118, 6)
            _add_round_box(x_label, y, label_w, row_h, lab, label_color, bold=True)
            _add_round_box(x_value, y, value_w, row_h, str(val), value_color, bold=False)

    # Table of selected long-form fields
    tbl_left = Cm(1.0)
    tbl_top = planned_tbl_top
    tbl_width = prs.slide_width - Cm(2.0)
    rows_count = sum(1 for _, key in DETAIL_FIELDS if str(row.get(key, "")).strip() != "")
    rows_count = max(rows_count, 1)
    table = slide.shapes.add_table(rows_count, 2, tbl_left, tbl_top, tbl_width, Cm(12)).table
    # Try to apply PowerPoint built-in table style "Medium Style 4 - Accent 1"
    # (French UI: "Moyen 4, accentuation 1"). Fallbacks if not available.
    for style_name in (
        "Medium Style 4 - Accent 1",
        "Medium Style 2 - Accent 1",
        "Medium Style 1 - Accent 1",
    ):
        try:
            table.style = style_name
            break
        except Exception:
            pass
    table.columns[0].width = Cm(6.0)
    table.columns[1].width = tbl_width - table.columns[0].width

    # Reduce inner paddings and line spacing to minimize extra vertical space
    for r_i in range(rows_count):
        for c_i in range(2):
            cell = table.cell(r_i, c_i)
            try:
                cell.margin_top = Cm(0.05)
                cell.margin_bottom = Cm(0.05)
            except Exception:
                pass
            tf = cell.text_frame
            for p in tf.paragraphs:
                try:
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    p.line_spacing = 1.0
                except Exception:
                    pass

    r = 0
    for label, key in DETAIL_FIELDS:
        val = str(row.get(key, "")).strip()
        if val == "":
            continue
        cell_lbl = table.cell(r, 0)
        cell_lbl.text = label
        # Bold label and set table font size to 8pt
        for p in cell_lbl.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(8)
                run.font.bold = True
        cell_val = table.cell(r, 1)
        cell_val.text = val
        for p in cell_val.text_frame.paragraphs:
            for run in p.runs:
                run.font.size = Pt(8)
        r += 1

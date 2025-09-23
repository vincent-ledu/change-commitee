from __future__ import annotations
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

from .utils import COLOR_MAP, DEFAULT_COLOR, hyperlink_for_rfc, set_title


def build_timeline_slide(prs: Presentation,
                         slide_index: int,
                         week_df: pd.DataFrame,
                         monday_next: datetime,
                         sunday_next: datetime,
                         title_text: str | None = None,
                         margins_cm=(1.0, 1.0, 5.0, 2.0),
                         col_gap_cm=0.15,
                         row_height_cm=1.0,
                         row_gap_cm=0.15) -> None:
    """
    Draw colored rounded rectangles for each change over a 7-day grid (S+1).
    - slide_index: index of the slide to draw on (typically 0 for template's first slide)
    - week_df: rows intersecting S+1, with parsed start_dt / end_dt
    """
    slide = prs.slides[slide_index]
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height

    left_cm, right_cm, top_cm, bottom_cm = margins_cm
    grid_left = Cm(left_cm)
    grid_top = Cm(top_cm)
    grid_width = SLIDE_WIDTH - Cm(left_cm + right_cm)
    grid_height = SLIDE_HEIGHT - Cm(top_cm + bottom_cm)

    col_count = 7
    col_gap = Cm(col_gap_cm)
    col_width = (grid_width - col_gap * (col_count - 1)) / col_count

    row_height = Cm(row_height_cm)
    # Box height tuned so up to 4 lines fit while keeping a slimmer footprint
    box_height = max(Cm(1.4), row_height)
    row_gap = Cm(row_gap_cm)

    MIN_BOX_WIDTH_CM = 5.0  # ensure room for "XXXXXXXX12345678" on first line
    row_segments: list[list[tuple[int, int]]] = []  # stores (left,right) per row in EMUs

    def day_index(dt):
        return max(0, min(6, (dt - monday_next).days))

    def clamp_to_week(dt: datetime) -> datetime:
        if dt < monday_next:
            return monday_next
        if dt > sunday_next:
            return sunday_next
        return dt

    def time_fraction_of_day(dt: datetime) -> float:
        # returns fraction in [0,1]
        seconds = dt.hour * 3600 + dt.minute * 60 + dt.second + dt.microsecond / 1_000_000.0
        return max(0.0, min(1.0, seconds / 86400.0))

    week_df = week_df.sort_values(["start_dt", "end_dt", "Numéro"], kind="stable")

    for _, row in week_df.iterrows():
        # Clamp to S+1 window for placement
        start_dt = clamp_to_week(row["start_dt"]) if isinstance(row["start_dt"], datetime) else row["start_dt"]
        end_dt = clamp_to_week(row["end_dt"]) if isinstance(row["end_dt"], datetime) else row["end_dt"]

        start_idx = day_index(start_dt)
        end_idx = day_index(end_dt)
        if end_idx < start_idx:
            end_idx = start_idx

        # Sub-day positioning using start/end hour within day
        start_frac = time_fraction_of_day(start_dt)
        end_frac = time_fraction_of_day(end_dt)

        # Compute left/right edges in slide coordinates
        left = grid_left + start_idx * (col_width + col_gap) + col_width * start_frac
        right = grid_left + end_idx * (col_width + col_gap) + col_width * end_frac
        # Enforce a minimum width sufficient to display "XXXXXXXX12345678"
        width = right - left
        if width < Cm(MIN_BOX_WIDTH_CM):
            width = Cm(MIN_BOX_WIDTH_CM)
            right = left + width
        # Detect whether we are exactly at the minimal enforced width
        is_min_width = int(width) == int(Cm(MIN_BOX_WIDTH_CM))
        left_emu = int(left)
        right_emu = int(right)

        # Find first row where the interval does not overlap existing boxes
        r_idx = None
        for idx, segments in enumerate(row_segments):
            overlap = any(not (right_emu <= seg_left or left_emu >= seg_right) for seg_left, seg_right in segments)
            if not overlap:
                r_idx = idx
                segments.append((left_emu, right_emu))
                break
        if r_idx is None:
            r_idx = len(row_segments)
            row_segments.append([(left_emu, right_emu)])

        top = grid_top + r_idx * (box_height + row_gap)
        height = box_height

        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

        typ_key = str(row.get("Type", "")).strip().lower()
        fill_color = COLOR_MAP.get(typ_key, DEFAULT_COLOR)
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill_color
        shp.line.width = Pt(0.75)
        shp.line.color.rgb = RGBColor(0, 0, 0)

        tf = shp.text_frame
        tf.clear()
        tf.word_wrap = True

        # If at minimal width (or otherwise very narrow), hide the resume but keep dates
        small_for_resume = is_min_width

        # First line: RFC number (hyperlink) and optional resume
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.LEFT
        rfc = str(row.get("Numéro", "")).strip()
        resume = str(row.get("Résumé", "")).strip()
        config_item = ""
        for key in ("Élément de configuration", "Element de configuration"):
            raw_val = row.get(key)
            if raw_val:
                config_item = str(raw_val).strip()
                break

        run_rfc = p1.add_run()
        run_rfc.text = rfc
        if rfc:
            run_rfc.hyperlink.address = hyperlink_for_rfc(rfc)
        run_rfc.font.bold = True
        run_rfc.font.size = Pt(8)
        run_rfc.font.color.rgb = RGBColor(255, 255, 255)

        if not small_for_resume and resume:
            run_sum = p1.add_run()
            run_sum.text = f" – {resume}"
            run_sum.font.size = Pt(8)
            run_sum.font.color.rgb = RGBColor(255, 255, 255)
        elif small_for_resume and config_item:
            # For narrow boxes, show the configuration element on the second line
            p_conf = tf.add_paragraph()
            p_conf.alignment = PP_ALIGN.LEFT
            run_conf = p_conf.add_run()
            run_conf.text = config_item
            run_conf.font.size = Pt(8)
            run_conf.font.color.rgb = RGBColor(255, 255, 255)

        # Dates must be kept even when the box is small for the resume
        # Next line: start planned date-time (third line when element of configuration is shown)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.LEFT
        try:
            start_str = row["start_dt"].strftime("%d/%m/%Y %H:%M")
        except Exception:
            start_str = str(row.get("Date de début planifiée", "")).strip()
        run_start = p2.add_run()
        run_start.text = start_str
        run_start.font.size = Pt(8)
        run_start.font.color.rgb = RGBColor(255, 255, 255)

        # Third line: end planned date-time
        p3 = tf.add_paragraph()
        p3.alignment = PP_ALIGN.LEFT
        try:
            end_str = row["end_dt"].strftime("%d/%m/%Y %H:%M")
        except Exception:
            end_str = str(row.get("Date de fin planifiée", "")).strip()
        run_end = p3.add_run()
        run_end.text = end_str
        run_end.font.size = Pt(8)
        run_end.font.color.rgb = RGBColor(255, 255, 255)

    if title_text:
        # Prefer title placeholder; fallback to adding a top textbox
        set_title(prs, slide, title_text)

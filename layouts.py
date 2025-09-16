from __future__ import annotations
from pptx import Presentation


def choose_detail_layout(prs: Presentation, layout_index: int | None):
    total = len(prs.slide_layouts)
    if total == 0:
        raise IndexError("Template has no slide layouts")

    # If user requested an index, use it safely (clamped) and warn if needed.
    if layout_index is not None:
        idx = max(0, min(layout_index, total - 1))
        if idx != layout_index:
            print(f"[WARN] detail layout index {layout_index} out of range; using {idx} instead (0..{total-1}).")
        return prs.slide_layouts[idx]

    # Auto-pick: try to find a 'blank' or 'vide' layout by name
    name_candidates = ("blank", "vide", "title only", "titre uniquement", "titre seul")
    for i, layout in enumerate(prs.slide_layouts):
        name = (layout.name or "").strip().lower()
        if any(k in name for k in name_candidates):
            return layout

    # Otherwise, pick the layout with the fewest placeholders (closest to blank)
    def placeholder_count(layout):
        try:
            return len(layout.placeholders)
        except Exception:
            return 999

    best = min(range(total), key=lambda i: placeholder_count(prs.slide_layouts[i]))
    return prs.slide_layouts[best]


def list_layouts(prs: Presentation):
    items = []
    for i, layout in enumerate(prs.slide_layouts):
        try:
            ph = len(layout.placeholders)
        except Exception:
            ph = "?"
        items.append((i, (layout.name or '').strip(), ph))
    return items


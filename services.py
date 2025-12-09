from __future__ import annotations
from datetime import datetime
import re
import unicodedata
import pandas as pd
from pptx import Presentation
from layouts import choose_detail_layout

from data_loader import load_dataset, try_guess_encoding, parse_fr_date
from periods import week_bounds_splus1, week_bounds_sminus1
from render.timeline import build_timeline_slide
from render.details import add_detail_slide
from render.charts import add_assignee_bar_chart_slide


REQUIRED_COLS = [
    "Numéro", "Type", "Etat", "Date de début planifiée", "Date de fin planifiée",
]


def _norm_label(s: str) -> str:
    s = str(s or "").strip().lower()
    # unifying quotes/apostrophes
    s = s.replace("’", "'").replace("`", "'")
    # normalize common separators to spaces
    s = s.replace("_", " ").replace("-", " ").replace("–", " ")
    # remove accents
    s = ''.join(c for c in unicodedata.normalize('NFD', s) if unicodedata.category(c) != 'Mn')
    # collapse whitespace
    s = re.sub(r"\s+", " ", s)
    # strip trailing punctuation like ':'
    s = s.strip(": ")
    return s


def _harmonize_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Rename incoming columns to canonical labels expected by the code, in an accent/quote-insensitive way.
    Only covers known fields used by the generator; unknown columns are left as-is.
    """
    canonical = [
        "Numéro", "Type", "Etat", "Date de début planifiée", "Date de fin planifiée",
        "Résumé", "Description", "Justification", "Plan d’implémentation",
        "Analyse des risques et de l’impact", "Plan de retour en arrière", "Plan de tests",
        "Informations complémentaires", "Groupe d’affectation", "Groupe gestionnaire",
        "Demandeur", "Affecté", "Affecté à", "Element de configuration", "CAB requis",
        "Code de fermeture", "Détail de clôture", "Balises",
    ]
    # Build normalized map for canonical labels; also add some common ASCII variants
    canon_map: dict[str, str] = {}
    for lab in canonical:
        canon_map.setdefault(_norm_label(lab), lab)
        # also provide an ASCII/straight-apostrophe variant mapping to the same
        canon_map.setdefault(_norm_label(lab.replace("’", "'")), lab)

    # Build rename mapping; avoid collisions if target already exists
    renames: dict[str, str] = {}
    existing_targets = set(df.columns)
    for col in df.columns:
        key = _norm_label(col)
        target = canon_map.get(key)
        if target and (col != target):
            # do not rename if it would overwrite an existing distinct column
            if target not in existing_targets or target == col:
                renames[col] = target
    if renames:
        df = df.rename(columns=renames)
    return df


def prepare_dataframe(data_path: str, encoding: str | None = None, sep: str | None = None) -> tuple[pd.DataFrame, dict]:
    df, meta = load_dataset(data_path, encoding=encoding, sep=sep)
    # Harmonize column labels (accent-insensitive, unify quotes)
    df = _harmonize_columns(df)
    missing = [c for c in REQUIRED_COLS if c not in df.columns]
    if missing:
        raise ValueError(f"Missing required column(s): {missing}. Present: {list(df.columns)}")

    # Parse dates robustly
    df["start_dt"] = df["Date de début planifiée"].apply(parse_fr_date)
    df["end_dt"] = df["Date de fin planifiée"].apply(parse_fr_date)
    return df, meta


def compute_weeks(ref_date: datetime) -> tuple[datetime, datetime, datetime, datetime]:
    _, monday_next, sunday_next = week_bounds_splus1(ref_date)
    monday_prev, sunday_prev = week_bounds_sminus1(ref_date)
    return monday_next, sunday_next, monday_prev, sunday_prev


def filter_week_df(df: pd.DataFrame, monday_next: datetime, sunday_next: datetime) -> pd.DataFrame:
    mask = (df["start_dt"] <= sunday_next) & (df["end_dt"] >= monday_next)
    return df.loc[mask].copy()


def filter_by_tags(df: pd.DataFrame, include_tags: list[str] | None, column: str = "Balises") -> pd.DataFrame:
    """Filter rows to keep only those whose `column` contains any of `include_tags`.
    - Matching is case-insensitive substring (robust to separators).
    - If include_tags is falsy or column is missing, returns df unchanged.
    """
    if not include_tags:
        return df
    if column not in df.columns:
        print(f"[WARN] Column '{column}' not found; --include-tags ignored.")
        return df
    tags = [t.strip() for t in include_tags if t and t.strip()]
    if not tags:
        return df
    pattern = "|".join(re.escape(t) for t in tags)
    mask = df[column].fillna("").astype(str).str.contains(pattern, case=False, regex=True)
    return df.loc[mask].copy()


def build_base_presentation(template_path: str,
                            week_df: pd.DataFrame,
                            monday_next: datetime,
                            sunday_next: datetime,
                            detail_layout_index: int | None = None,
                            splus1_layout_index: int | None = None,
                            assignee_layout_index: int | None = None) -> Presentation:
    prs = Presentation(template_path)
    # Place the S+1 timeline either on the template's first slide (default)
    # or on a newly added slide using the requested layout index.
    # Prepare per-type dataframes
    def _norm(s: str) -> str:
        import unicodedata
        s = str(s or '').strip().lower()
        s = ''.join(c for c in unicodedata.normalize('NFD', s) if unicodedata.category(c) != 'Mn')
        return s
    type_series = week_df['Type'].fillna('').astype(str).map(_norm)
    info_labels = {'pour information', 'information', 'info'}
    df_urgent = week_df.loc[type_series == 'urgent']
    df_normal = week_df.loc[type_series == 'normal']
    df_agile = week_df.loc[type_series == 'agile']
    df_info = week_df.loc[type_series.isin(info_labels)]

    # Decide where to render: use template's first slide for the first timeline if no layout index
    if splus1_layout_index is None:
        # Slide 0: Urgent
        title_u = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Urgent"
        build_timeline_slide(prs, slide_index=0, week_df=df_urgent,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_u)
        # Add slide: Normal
        layout = choose_detail_layout(prs, layout_index=None)
        prs.slides.add_slide(layout)
        title_n = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Normal"
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_normal,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_n)
        # Add slide: Agile
        layout = choose_detail_layout(prs, layout_index=None)
        prs.slides.add_slide(layout)
        title_a = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Agile"
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_agile,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_a)
        layout = choose_detail_layout(prs, layout_index=None)
        prs.slides.add_slide(layout)
        title_i = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Pour information"
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_info,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_i)
    else:
        # Use the chosen layout for all slides (Urgent → Normal → Agile → Pour information)
        layout = choose_detail_layout(prs, layout_index=splus1_layout_index)
        prs.slides.add_slide(layout)
        title_u = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Urgent"
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_urgent,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_u)
        prs.slides.add_slide(layout)
        title_n = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Normal"
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_normal,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_n)
        prs.slides.add_slide(layout)
        title_a = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Agile"
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_agile,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_a)
        prs.slides.add_slide(layout)
        title_i = f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — Pour information"
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_info,
                             monday_next=monday_next, sunday_next=sunday_next,
                             title_text=title_i)
    chart_title = (
        f"Changements S+1 ({monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}) — répartition par affecté"
    )
    chart_layout_index = assignee_layout_index if assignee_layout_index is not None else splus1_layout_index
    add_assignee_bar_chart_slide(
        prs,
        week_df,
        chart_title,
        layout_index=chart_layout_index,
    )
    def _type_priority(norm_label: str) -> int:
        if norm_label == 'urgent':
            return 0
        if norm_label == 'normal':
            return 1
        if norm_label == 'agile':
            return 2
        if norm_label in info_labels:
            return 3
        return 4

    week_df = week_df.assign(_type_order=type_series.map(_type_priority))
    week_df = week_df.sort_values(["_type_order", "start_dt", "Numéro"], kind="stable")
    week_df = week_df.drop(columns="_type_order")
    for _, row in week_df.iterrows():
        add_detail_slide(prs, row, layout_index=detail_layout_index)
    return prs

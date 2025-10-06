#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Generate a Change Advisory Board (CAB) support PowerPoint:

What it builds
- S+1 timeline: boxes positioned from start to end date (with intra-day hour precision), colored by Type.
- One detail slide per change included in S+1 (RFC + fields table).
- Optional S-1 slides: a pie chart by closure code, and a table listing non-“Succès”.
- Optional current week (S) timeline slide.

Visuals
- Colors: Urgent → orange, Normal → blue, Agile → green.
- RFC number is clickable: https://outils.change.fr/change=rfc123
- Timeline box text:
  - Line 1: “RFC – Résumé” (Résumé omitted on very narrow boxes).
  - Lines 2–3: start and end planned date-times (always displayed), with smaller font.

CLI (most common)
  python generate_cab_pptx.py \
    --data cab_changes.xlsx \
    --template template_change_SPLUS1.pptx \
    --out change_generated.pptx \
    --ref-date 2025-09-09 \
    --sminus1-pie --current-week

Main options
- --data PATH                      CSV/Excel containing changes (required)
- --template PATH                  PPTX template (required)
- --out PATH                       Output PPTX (required)
- --ref-date YYYY-MM-DD            Reference date; defaults to today
- --detail-layout-index N          Layout index for detail slides
- --splus1-layout-index N          Layout index for S+1 timeline (otherwise uses template’s first slide)
- --sminus1-pie                    Add S-1 statistics slides (pie + non‑“Succès” list)
- --sminus1-layout-index N         Layout index for S-1 slides
- --current-week                   Add current week (S) timeline at the end
- --current-week-layout-index N    Layout index for the current week slide
- --list-layouts                   Print available template layouts and exit
- --encoding ENC                   Force CSV encoding (e.g. cp1252, latin1, utf-8-sig)
- --sep SEP                        Force CSV separator (e.g. ';' ',' '\t')

Data requirements
- Required columns: “Numéro”, “Type”, “Etat”, “Date de début planifiée”, “Date de fin planifiée”.
- Dates are parsed robustly (dd/mm/yy, dd/mm/YYYY, YYYY-MM-DD, with/without HH:MM[:SS]).

Dependencies
- pip install python-pptx pandas python-dateutil openpyxl
- Optional (auto-encoding for CSV): pip install charset-normalizer
"""
from __future__ import annotations
import argparse
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
# Extracted modules (step 1 of refactor)
from layouts import choose_detail_layout, list_layouts
from render.utils import hyperlink_for_rfc, set_rfc_url_template, set_timeline_color_map
from services import prepare_dataframe, compute_weeks, filter_week_df, build_base_presentation, filter_by_tags
from render.timeline import build_timeline_slide
from periods import week_bounds_current
from config import load_config

# ---------------------- Helpers ----------------------

def parse_fr_date(s) -> datetime:
    """Parse une date ou date+heure en acceptant plusieurs formats.
    Gère aussi directement les objets datetime/pandas.Timestamp.
    """
    # Déjà un datetime ?
    if isinstance(s, datetime):
        return s
    try:
        import pandas as _pd
        if isinstance(s, _pd.Timestamp):
            return s.to_pydatetime()
    except Exception:
        pass

    s = str(s or "").strip()
    # Essais explicites communs (avec et sans heure)
    fmts = (
        "%d/%m/%y",
        "%d/%m/%Y",
        "%Y-%m-%d",
        "%d/%m/%y %H:%M",
        "%d/%m/%Y %H:%M",
        "%d/%m/%y %H:%M:%S",
        "%d/%m/%Y %H:%M:%S",
        "%Y-%m-%d %H:%M",
        "%Y-%m-%d %H:%M:%S",
        "%Y-%m-%dT%H:%M",
        "%Y-%m-%dT%H:%M:%S",
    )
    for fmt in fmts:
        try:
            return datetime.strptime(s, fmt)
        except Exception:
            pass

    # Dernier recours: pandas.to_datetime avec dayfirst
    try:
        dt = pd.to_datetime(s, dayfirst=True, errors="raise")
        return dt.to_pydatetime() if hasattr(dt, "to_pydatetime") else dt
    except Exception:
        pass

    raise ValueError(f"Unrecognized date/datetime format: {s!r}")

# week bounds now provided by periods.py

# Colors for S-1 closure code pie chart
PIE_COLOR_SUCCESS = RGBColor(0, 176, 80)            # green
PIE_COLOR_SUCCESS_DIFFICULT = RGBColor(255, 192, 0) # yellow
PIE_COLOR_PARTIAL = RGBColor(255, 204, 153)         # light orange
PIE_COLOR_FAIL_NO_ROLLBACK = RGBColor(255, 140, 0)  # dark orange
PIE_COLOR_FAIL_ROLLBACK = RGBColor(192, 0, 0)       # red

def _norm_text(s: str) -> str:
    import unicodedata
    s = str(s or "").strip().lower()
    s = ''.join(c for c in unicodedata.normalize('NFD', s) if unicodedata.category(c) != 'Mn')
    return s

def classify_closure(code: str) -> str:
    """
    Classify 'Code de fermeture' into one of the buckets:
    - succes
    - succes_difficulte
    - partiel
    - echec_sans_retour
    - echec_avec_retour
    Unknown or empty returns '' (ignored).
    """
    t = _norm_text(code)
    if not t:
        return ''
    # direct keywords
    # success bucket (accent-insensitive)
    if 'succes' in t or 'reussi' in t:
        if 'diffic' in t:
            return 'succes_difficulte'
        return 'succes'
    if 'partiel' in t or 'partial' in t:
        return 'partiel'
    if 'echec' in t or 'fail' in t:
        if 'retour' in t or 'rollback' in t:
            if 'sans' in t and ('retour' in t or 'rollback' in t):
                # e.g., "echec sans retour arriere"
                return 'echec_sans_retour'
            return 'echec_avec_retour'
        # no explicit mention of rollback -> assume without rollback
        return 'echec_sans_retour'
    # explicit rollback mention without failure keyword -> consider as with rollback
    if 'rollback' in t or 'retour arriere' in t:
        return 'echec_avec_retour'
    return ''

def add_sminus1_pie_slide(prs: Presentation,
                          df: pd.DataFrame,
                          monday_prev: datetime,
                          sunday_prev: datetime,
                          layout_index: int | None = None,
                          closure_col: str = 'Code de fermeture') -> None:
    """
    Add a pie chart slide summarizing S-1 by closure code categories.
    Selection criterion: rows with end date within [monday_prev, sunday_prev].
    """
    # Choose layout
    layout = choose_detail_layout(prs, layout_index)
    slide = prs.slides.add_slide(layout)

    # Title (use placeholder if available)
    from render.utils import set_title
    set_title(prs, slide, f"Changements S-1 par code de fermeture ({monday_prev.strftime('%d/%m/%Y')} → {sunday_prev.strftime('%d/%m/%Y')})")

    # Aggregate counts
    subset = df.copy()
    mask_end_in_week = (subset['end_dt'] >= monday_prev) & (subset['end_dt'] <= sunday_prev)
    subset = subset.loc[mask_end_in_week]

    buckets = {
        'succes': 0,
        'succes_difficulte': 0,
        'partiel': 0,
        'echec_sans_retour': 0,
        'echec_avec_retour': 0,
    }
    if closure_col in subset.columns:
        for val in subset[closure_col].fillna(''):
            key = classify_closure(val)
            if key in buckets:
                buckets[key] += 1

    # Prepare chart data in desired order
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE

    categories = [
        'Succès',
        'Succès avec difficulté',
        'Implémenté partiellement',
        'Échec sans retour arrière',
        'Échec avec retour arrière',
    ]
    keys_order = ['succes', 'succes_difficulte', 'partiel', 'echec_sans_retour', 'echec_avec_retour']
    values = [buckets[k] for k in keys_order]

    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series('Changements', values)

    # Add chart
    chart_left = Cm(3.0)
    chart_top = Cm(3.0)
    chart_width = Cm(20.0)
    chart_height = Cm(12.0)
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, chart_left, chart_top, chart_width, chart_height, chart_data
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.position = 2  # right
    chart.has_title = False
    chart.plots[0].has_data_labels = True
    dl = chart.plots[0].data_labels
    dl.show_category_name = True
    dl.show_percentage = True

    # Apply colors per slice
    colors = [
        PIE_COLOR_SUCCESS,
        PIE_COLOR_SUCCESS_DIFFICULT,
        PIE_COLOR_PARTIAL,
        PIE_COLOR_FAIL_NO_ROLLBACK,
        PIE_COLOR_FAIL_ROLLBACK,
    ]
    series = chart.series[0]
    for idx, point in enumerate(series.points):
        fill = point.format.fill
        fill.solid()
        fill.fore_color.rgb = colors[idx]

def _first_present_col(df: pd.DataFrame, candidates: list[str]) -> str | None:
    for c in candidates:
        if c in df.columns:
            return c
    return None

def add_sminus1_non_success_slide(prs: Presentation,
                                   df: pd.DataFrame,
                                   monday_prev: datetime,
                                   sunday_prev: datetime,
                                   layout_index: int | None = None,
                                   closure_col: str = 'Code de fermeture') -> None:
    """Add a table slide listing S-1 changes that are not 'Succès'."""
    layout = choose_detail_layout(prs, layout_index)
    slide = prs.slides.add_slide(layout)

    # Title (use placeholder if available)
    from render.utils import set_title
    set_title(prs, slide, f"Changements S-1 non 'Succès' ({monday_prev.strftime('%d/%m/%Y')} → {sunday_prev.strftime('%d/%m/%Y')})")

    # Determine columns
    resume_col = _first_present_col(df, ['Résumé', 'Resume'])
    state_col = _first_present_col(df, ['État', 'Etat', 'Etat du changement', 'Etat du changements'])
    detail_col = _first_present_col(df, ['Détail de clôture', 'Detail de clôture', 'Détail de cloture', 'Detail de cloture'])

    subset = df.copy()
    mask_end_in_week = (subset['end_dt'] >= monday_prev) & (subset['end_dt'] <= sunday_prev)
    subset = subset.loc[mask_end_in_week]

    # Filter: not strictly 'succes'
    rows = []
    for _, r in subset.iterrows():
        code_val = r.get(closure_col, '')
        key = classify_closure(code_val)
        if key == 'succes':
            continue
        rows.append(r)

    if not rows:
        # No data: show a friendly message
        msg_width = prs.slide_width - Cm(2.0)
        msg_shape = slide.shapes.add_textbox(Cm(1.0), Cm(3.0), msg_width, Cm(3.0))
        msg_tf = msg_shape.text_frame
        msg_tf.text = "Aucun changement non 'Succès' pour S-1."
        return

    # Build table with header + items
    tbl_left = Cm(1.0); tbl_top = Cm(3.0)
    tbl_width = prs.slide_width - Cm(2.0)
    headers = ['Numéro', 'Résumé', 'Etat', 'Code de fermeture', 'Détail de clôture']
    table = slide.shapes.add_table(len(rows) + 1, len(headers), tbl_left, tbl_top, tbl_width, Cm(12)).table

    def _set_cell_font(cell, size_pt: float | None = None, bold: bool | None = None) -> None:
        tf = cell.text_frame
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                if size_pt is not None:
                    run.font.size = Pt(size_pt)
                if bold is not None:
                    run.font.bold = bold

    # Set column widths (approximate for readability)
    col_widths = [Cm(4.0), Cm(9.0), Cm(4.0), Cm(5.0), tbl_width - Cm(4.0 + 9.0 + 4.0 + 5.0)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        _set_cell_font(cell, bold=True)

    # Data rows
    r_idx = 1
    for r in rows:
        # Numéro with hyperlink
        cell_num = table.cell(r_idx, 0)
        tf = cell_num.text_frame
        tf.clear()
        p0 = tf.paragraphs[0]
        run0 = p0.add_run()
        rfc = str(r.get('Numéro', '')).strip()
        run0.text = rfc
        if rfc:
            run0.hyperlink.address = hyperlink_for_rfc(rfc)
            run0.font.bold = True
        # Résumé
        cell_res = table.cell(r_idx, 1)
        cell_res.text = str(r.get(resume_col, '')) if resume_col else ''

        # Etat
        cell_state = table.cell(r_idx, 2)
        cell_state.text = str(r.get(state_col, '')) if state_col else ''

        # Code de fermeture
        cell_code = table.cell(r_idx, 3)
        cell_code.text = str(r.get(closure_col, ''))

        # Détail de clôture
        cell_det = table.cell(r_idx, 4)
        cell_det.text = str(r.get(detail_col, '')) if detail_col else ''

        r_idx += 1

# dataset loader now provided by data_loader.py

# timeline and detail slide builders now provided by render.timeline and render.details

# ---------------------- Main ----------------------

def main():
    ap = argparse.ArgumentParser(description="Generate CAB PowerPoint (S+1 timeline + details).")
    ap.add_argument("--data", required=True, help="Path to CSV/Excel with changes")
    ap.add_argument("--config", default=None, help="Path to JSON config file for default options")
    ap.add_argument("--template", required=False, help="Path to PPTX template (timeline on slide 0)")
    ap.add_argument("--out", required=False, help="Output PPTX path")
    ap.add_argument("--ref-date", default=None, help="Reference date (YYYY-MM-DD); default: today")
    ap.add_argument("--detail-layout-index", type=int, default=None, help="Optional slide layout index to use for detail slides")
    ap.add_argument("--sminus1-pie", action="store_true", help="Add an S-1 pie chart slide by closure code")
    ap.add_argument("--sminus1-layout-index", type=int, default=None, help="Optional slide layout index for S-1 pie slide")
    ap.add_argument("--list-layouts", action="store_true", help="List slide layouts in the template and exit")
    ap.add_argument("--current-week", action="store_true", help="Add a timeline slide for the current week (S)")
    ap.add_argument("--current-week-layout-index", type=int, default=None, help="Optional slide layout index for the current week slide")
    ap.add_argument("--encoding", default=None, help="Force CSV encoding (e.g. cp1252, latin1, utf-8-sig)")
    ap.add_argument("--sep", default=None, help="Force CSV separator (e.g. ';' ',' '\\t'). If omitted, auto-try common ones.")
    ap.add_argument("--include-tags", default=None, help="Comma-separated tags to include (matches column 'Balises'). Example: RED_TRUC-TEL,GRE_BIDULE-PDT")
    ap.add_argument("--splus1-layout-index", type=int, default=None, help="Optional slide layout index for S+1 timeline slide (otherwise uses template's first slide)")
    args = ap.parse_args()

    # Load config file (JSON) and use it to fill missing options; CLI overrides config
    cfg = load_config(args.config) if getattr(args, 'config', None) else {}
    def _cfg(name, default=None):
        return cfg.get(name, default)

    # Branding / visuals overrides
    set_rfc_url_template(_cfg('rfc_base_url') or _cfg('rfc_url_template'))
    set_timeline_color_map(_cfg('timeline_colors'))

    # Merge config defaults where CLI is unset/False
    if not getattr(args, 'template', None):
        setattr(args, 'template', _cfg('template', None))
    if not getattr(args, 'out', None):
        setattr(args, 'out', _cfg('out', None))
    if not getattr(args, 'ref_date', None) and _cfg('ref_date'):
        setattr(args, 'ref_date', _cfg('ref_date'))
    if getattr(args, 'detail_layout_index', None) is None and _cfg('detail_layout_index') is not None:
        setattr(args, 'detail_layout_index', int(_cfg('detail_layout_index')))
    if not getattr(args, 'sminus1_pie', False) and bool(_cfg('sminus1_pie', False)):
        setattr(args, 'sminus1_pie', True)
    if getattr(args, 'sminus1_layout_index', None) is None and _cfg('sminus1_layout_index') is not None:
        setattr(args, 'sminus1_layout_index', int(_cfg('sminus1_layout_index')))
    if not getattr(args, 'list_layouts', False) and bool(_cfg('list_layouts', False)):
        setattr(args, 'list_layouts', True)
    if not getattr(args, 'encoding', None) and _cfg('encoding'):
        setattr(args, 'encoding', _cfg('encoding'))
    if getattr(args, 'sep', None) is None and (_cfg('sep') is not None):
        setattr(args, 'sep', _cfg('sep'))
    if getattr(args, 'splus1_layout_index', None) is None and _cfg('splus1_layout_index') is not None:
        setattr(args, 'splus1_layout_index', int(_cfg('splus1_layout_index')))
    if not getattr(args, 'current_week', False) and bool(_cfg('current_week', False)):
        setattr(args, 'current_week', True)
    if getattr(args, 'current_week_layout_index', None) is None and _cfg('current_week_layout_index') is not None:
        setattr(args, 'current_week_layout_index', int(_cfg('current_week_layout_index')))
    if not getattr(args, 'include_tags', None) and _cfg('include_tags'):
        tags = _cfg('include_tags')
        if isinstance(tags, list):
            setattr(args, 'include_tags', ",".join(str(t) for t in tags))
        else:
            setattr(args, 'include_tags', str(tags))

    # Validate required values possibly supplied via config
    if not getattr(args, 'template', None):
        raise SystemExit("error: --template is required (can be provided via --config)")
    if not getattr(args, 'out', None):
        raise SystemExit("error: --out is required (can be provided via --config)")

    df, meta = prepare_dataframe(args.data, encoding=args.encoding, sep=args.sep)
    # Optional: filter by tags from 'Balises' column
    if args.include_tags:
        tags = [t.strip() for t in str(args.include_tags).split(',') if t.strip()]
        before = len(df)
        df = filter_by_tags(df, tags, column="Balises")
        print(f"[INFO] Tag filter applied ({len(tags)} tag(s)): {before} → {len(df)} rows")
    print(f"[INFO] Loaded dataset via {meta.get('reader')} "
          f"(encoding={meta.get('encoding')}, sep={repr(meta.get('sep'))}, engine={meta.get('engine')})")

    if args.ref_date:
        ref_date = datetime.strptime(args.ref_date, "%Y-%m-%d")
    else:
        ref_date = datetime.today()

    monday_next, sunday_next, monday_prev, sunday_prev = compute_weeks(ref_date)
    week_df = filter_week_df(df, monday_next, sunday_next)
    print(f"[INFO] Changes in S+1: {len(week_df)}")

    # If only listing layouts, do it and exit early
    if args.list_layouts:
        prs = Presentation(args.template)
        print("[INFO] Available slide layouts (index: name | placeholders)")
        for i, name, ph in list_layouts(prs):
            print(f"  {i}: {name} | placeholders={ph}")
        return
    prs = build_base_presentation(
        template_path=args.template,
        week_df=week_df,
        monday_next=monday_next,
        sunday_next=sunday_next,
        detail_layout_index=args.detail_layout_index,
        splus1_layout_index=args.splus1_layout_index,
    )

    # Then S-1 statistics and non-success list
    if args.sminus1_pie:
        add_sminus1_pie_slide(prs, df=df, monday_prev=monday_prev, sunday_prev=sunday_prev,
                              layout_index=args.sminus1_layout_index)
        add_sminus1_non_success_slide(prs, df=df, monday_prev=monday_prev, sunday_prev=sunday_prev,
                                      layout_index=args.sminus1_layout_index)

    # Finally, optionally add a timeline slide for the current week (S)
    if args.current_week:
        monday_cur, sunday_cur = week_bounds_current(ref_date)
        curr_week_df = filter_week_df(df, monday_cur, sunday_cur)
        # Split by Type (normal/agile) using accent-insensitive normalization
        import unicodedata
        def _norm(s: str) -> str:
            s = str(s or '').strip().lower()
            return ''.join(c for c in unicodedata.normalize('NFD', s) if unicodedata.category(c) != 'Mn')
        type_series = curr_week_df['Type'].fillna('').astype(str).map(_norm)
        df_norm = curr_week_df.loc[type_series == 'normal']
        df_agil = curr_week_df.loc[type_series == 'agile']

        # Add slide: Normal
        layout = choose_detail_layout(prs, layout_index=args.current_week_layout_index)
        prs.slides.add_slide(layout)
        from render.utils import set_title
        set_title(prs, prs.slides[-1], f"Changements cette semaine ({monday_cur.strftime('%d/%m/%Y')} → {sunday_cur.strftime('%d/%m/%Y')}) — Normal")
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_norm,
                             monday_next=monday_cur, sunday_next=sunday_cur)

        # Add slide: Agile
        layout = choose_detail_layout(prs, layout_index=args.current_week_layout_index)
        prs.slides.add_slide(layout)
        set_title(prs, prs.slides[-1], f"Changements cette semaine ({monday_cur.strftime('%d/%m/%Y')} → {sunday_cur.strftime('%d/%m/%Y')}) — Agile")
        build_timeline_slide(prs, slide_index=len(prs.slides) - 1, week_df=df_agil,
                             monday_next=monday_cur, sunday_next=sunday_cur)

    prs.save(args.out)
    print(f"[OK] Generated: {args.out}")
    print(f"[OK] S+1 week: {monday_next.strftime('%d/%m/%Y')} → {sunday_next.strftime('%d/%m/%Y')}")
    if args.sminus1_pie:
        print(f"[OK] S-1 week: {monday_prev.strftime('%d/%m/%Y')} → {sunday_prev.strftime('%d/%m/%Y')}")

if __name__ == "__main__":
    main()
